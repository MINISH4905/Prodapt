from pptx import Presentation
from pptx.util import Pt
import os


def add_slide(prs, title, content):

    slide = prs.slides.add_slide(
        prs.slide_layouts[1]
    )

    slide.shapes.title.text = title

    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()

    if isinstance(content, list):

        for index, item in enumerate(content):

            if index == 0:
                paragraph = text_frame.paragraphs[0]

            else:
                paragraph = text_frame.add_paragraph()

            paragraph.text = str(item)
            paragraph.font.size = Pt(20)

    else:

        text_frame.paragraphs[0].text = str(content)
        text_frame.paragraphs[0].font.size = Pt(20)

    return slide


def generate_pitch(data, output_path):

    os.makedirs(os.path.dirname(output_path), exist_ok=True)

    prs = Presentation()

    startup = data.startup
    refined = data.refined_idea

    startup_name = startup.get(
        "name",
        "Startup"
    )

    # --------------------------------
    # SLIDE 1
    # --------------------------------

    slide = prs.slides.add_slide(
        prs.slide_layouts[0]
    )

    slide.shapes.title.text = startup_name

    slide.placeholders[1].text = (
        "VentureX-Ray Final Investor Pitch"
    )

    # --------------------------------
    # SLIDE 2
    # --------------------------------

    add_slide(
        prs,
        "The Problem",
        startup.get(
            "problem",
            "Problem statement not provided."
        )
    )

    # --------------------------------
    # SLIDE 3
    # --------------------------------

    add_slide(
        prs,
        "Our Solution",
        refined.get(
            "solution",
            refined.get(
                "description",
                "Solution not provided."
            )
        )
    )

    # --------------------------------
    # SLIDE 4
    # --------------------------------

    add_slide(
        prs,
        "Target Market",
        startup.get(
            "target_market",
            "Target market information not provided."
        )
    )

    # --------------------------------
    # SLIDE 5
    # --------------------------------

    add_slide(
        prs,
        "Business Model",
        startup.get(
            "business_model",
            "Business model information not provided."
        )
    )

    # --------------------------------
    # SLIDE 6
    # --------------------------------

    attacker_content = []

    for index, attacker in enumerate(
        data.attacker_results,
        start=1
    ):

        if isinstance(attacker, dict):

            text = ", ".join(
                f"{key}: {value}"
                for key, value in attacker.items()
            )

        else:

            text = str(attacker)

        attacker_content.append(
            f"Attacker {index}: {text}"
        )

    if not attacker_content:

        attacker_content = [
            "No attacker findings available."
        ]

    add_slide(
        prs,
        "AI Stress-Test Findings",
        attacker_content
    )

    # --------------------------------
    # SLIDE 7
    # --------------------------------

    vulnerability_content = []

    for vulnerability in data.vulnerabilities:

        if isinstance(vulnerability, dict):

            text = ", ".join(
                f"{key}: {value}"
                for key, value in vulnerability.items()
            )

        else:

            text = str(vulnerability)

        vulnerability_content.append(text)

    if not vulnerability_content:

        vulnerability_content = [
            "No vulnerabilities available."
        ]

    add_slide(
        prs,
        "Key Vulnerabilities",
        vulnerability_content
    )

    # --------------------------------
    # SLIDE 8
    # --------------------------------

    add_slide(
        prs,
        "Founder Defense",
        [
            f"Defense Score: {data.defense_score}/100",
            "Founder responses were evaluated against AI attacks."
        ]
    )

    # --------------------------------
    # SLIDE 9
    # --------------------------------

    investor_content = []

    for conversation in data.investor_conversation:

        if isinstance(conversation, dict):

            text = ", ".join(
                f"{key}: {value}"
                for key, value in conversation.items()
            )

        else:

            text = str(conversation)

        investor_content.append(text)

    if not investor_content:

        investor_content = [
            "No investor simulation data available."
        ]

    add_slide(
        prs,
        "Investor Simulation",
        investor_content
    )

    # --------------------------------
    # SLIDE 10
    # --------------------------------

    add_slide(
        prs,
        "Investment Readiness Score",
        [
            f"{data.defense_score}/100",
            f"Final Decision: {data.decision}"
        ]
    )

    # --------------------------------
    # SLIDE 11
    # --------------------------------

    if data.decision.upper() == "STRONG":

        recommendation = [
            "Startup is investment ready.",
            "Major vulnerabilities were addressed.",
            "Continue validating the business model."
        ]

    else:

        recommendation = [
            "Startup requires further refinement.",
            "Address major vulnerabilities.",
            "Repeat the stress-testing process."
        ]

    add_slide(
        prs,
        "Final Recommendation",
        recommendation
    )

    # --------------------------------
    # SLIDE 12
    # --------------------------------

    add_slide(
        prs,
        "Final Pitch",
        [
            startup_name,
            "Stress-tested. Defended. Investor-ready."
        ]
    )

    prs.save(output_path)

    return output_path